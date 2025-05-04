import streamlit as st
import requests
import openai
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation

# --- API KEYS ---
SERP_API_KEY = "cb77062e09a413a5156e5a92440ab298d05e33ddc00a551a852b240fca149f22"
NEWS_API_KEY = "7271cff2fd6949f4b617adfa2f44d6d5"
openai.api_key = "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"

# --- Data Functions ---
def get_company_overview(company_name):
    params = {
        "engine": "google",
        "q": f"{company_name} company overview",
        "api_key": SERP_API_KEY
    }
    url = "https://serpapi.com/search"
    response = requests.get(url, params=params)
    data = response.json()
    return {
        "history": data.get("knowledge_graph", {}).get("description", ""),
        "mission": data.get("knowledge_graph", {}).get("mission", ""),
        "leadership": data.get("knowledge_graph", {}).get("founders", []),
        "locations": data.get("knowledge_graph", {}).get("headquarters", "")
    }

def get_financial_data(company_name, time_period):
    params = {
        "engine": "google_finance",
        "q": company_name,
        "api_key": SERP_API_KEY
    }
    url = "https://serpapi.com/search"
    response = requests.get(url, params=params)
    data = response.json()
    # Example structure, adapt as needed
    return {
        "annual_revenue": [
            {"year": "2020", "value": 1000},
            {"year": "2021", "value": 1200},
            {"year": "2022", "value": 1500},
        ],
        "employee_count": data.get("financial_data", {}).get("employees", 0),
        "key_ratios": data.get("financial_data", {}).get("ratios", {}),
        "statements": {
            "income": data.get("financial_data", {}).get("income_statement", {}),
            "balance": data.get("financial_data", {}).get("balance_sheet", {}),
            "cash_flow": data.get("financial_data", {}).get("cash_flow_statement", {})
        }
    }

def get_news(company_name, time_period):
    url = "https://newsapi.org/v2/everything"
    params = {
        "q": company_name,
        "from": time_period,
        "sortBy": "relevancy",
        "apiKey": NEWS_API_KEY,
        "language": "en"
    }
    response = requests.get(url, params=params)
    return response.json().get("articles", [])

def generate_swot(overview, financials, news):
    prompt = f"""Generate a detailed SWOT analysis for the following company:
Overview: {overview}
Financials: {financials}
Recent News: {news}
"""
    client = openai.OpenAI(api_key=openai.api_key)
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def generate_benchmark(company_name, industry_sector):
    prompt = f"Provide a benchmark analysis for {company_name} in the {industry_sector} sector."
    client = openai.OpenAI(api_key=openai.api_key)
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

def generate_industry_intel(industry_sector, geographic_focus):
    prompt = f"Provide industry intelligence for the {industry_sector} sector in {geographic_focus}."
    client = openai.OpenAI(api_key=openai.api_key)
    response = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content

# --- Report Generators ---
def create_revenue_chart(revenue_data):
    years = [item['year'] for item in revenue_data]
    values = [item['value'] for item in revenue_data]
    fig, ax = plt.subplots()
    ax.plot(years, values, marker='o')
    ax.set_title('Annual Revenue Trend')
    ax.set_xlabel('Year')
    ax.set_ylabel('Revenue')
    buf = BytesIO()
    plt.savefig(buf, format='png')
    plt.close(fig)
    buf.seek(0)
    return buf

def generate_pdf_report(overview, financials, swot, benchmark, industry_intel):
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    c.setFont("Helvetica-Bold", 16)
    c.drawString(30, height - 50, "Company Analysis Report")
    c.setFont("Helvetica", 12)
    c.drawString(30, height - 80, "Executive Summary")
    c.drawString(30, height - 100, f"Overview: {overview}")
    c.drawString(30, height - 120, f"Financials: {financials}")
    c.drawString(30, height - 140, f"SWOT: {swot}")
    c.drawString(30, height - 160, f"Benchmark: {benchmark}")
    c.drawString(30, height - 180, f"Industry Intelligence: {industry_intel}")
    c.save()
    buf.seek(0)
    return buf

def generate_excel_report(overview, financials, swot, benchmark, industry_intel):
    output = BytesIO()
    with pd.ExcelWriter(output, engine='xlsxwriter') as writer:
        pd.DataFrame([overview]).to_excel(writer, sheet_name='Overview')
        pd.DataFrame(financials['annual_revenue']).to_excel(writer, sheet_name='Revenue')
        pd.DataFrame([{"SWOT": swot}]).to_excel(writer, sheet_name='SWOT')
        pd.DataFrame([{"Benchmark": benchmark}]).to_excel(writer, sheet_name='Benchmark')
        pd.DataFrame([{"Industry Intelligence": industry_intel}]).to_excel(writer, sheet_name='IndustryIntel')
    output.seek(0)
    return output

def generate_ppt_report(overview, financials, swot, benchmark, industry_intel):
    output = BytesIO()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Company Analysis Report"
    slide.placeholders[1].text = f"Overview: {overview['history']}"
    prs.save(output)
    output.seek(0)
    return output

# --- Streamlit UI ---
st.title("Market Research Automation Tool")

with st.form("input_form"):
    company_name = st.text_input("Company Name")
    industry_sector = st.text_input("Industry Sector")
    time_period = st.text_input("Time Period (YYYY-MM-DD)", value="2022-01-01")
    geographic_focus = st.text_input("Geographic Focus", value="Global")
    submitted = st.form_submit_button("Generate Report")

if submitted:
    with st.spinner("Collecting data and generating report..."):
        overview = get_company_overview(company_name)
        financials = get_financial_data(company_name, time_period)
        news = get_news(company_name, time_period)
        swot = generate_swot(overview, financials, news)
        benchmark = generate_benchmark(company_name, industry_sector)
        industry_intel = generate_industry_intel(industry_sector, geographic_focus)

        st.subheader("Company Overview")
        st.json(overview)
        st.subheader("Financials")
        st.json(financials)
        st.subheader("SWOT Analysis")
        st.write(swot)
        st.subheader("Benchmark Analysis")
        st.write(benchmark)
        st.subheader("Industry Intelligence")
        st.write(industry_intel)

        # Chart
        st.subheader("Annual Revenue Trend")
        chart_buf = create_revenue_chart(financials['annual_revenue'])
        st.image(chart_buf)

        # Download buttons
        pdf_buf = generate_pdf_report(overview, financials, swot, benchmark, industry_intel)
        st.download_button("Download PDF Report", pdf_buf, file_name="company_report.pdf", mime="application/pdf")

        excel_buf = generate_excel_report(overview, financials, swot, benchmark, industry_intel)
        st.download_button("Download Excel Report", excel_buf, file_name="company_report.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

        ppt_buf = generate_ppt_report(overview, financials, swot, benchmark, industry_intel)
        st.download_button("Download PowerPoint Report", ppt_buf, file_name="company_report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
